"""Load the Northwestern template, strip its sample slides, and (later) render
an outline into a finished deck."""
import sys
import json
import traceback
from pptx import Presentation
import nwstyle as ns
import components

_R_ID = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"


def load_template(path):
    return Presentation(path)


def clear_slides(prs):
    """Remove every slide AND its relationship so re-adding slides does not
    collide with orphaned parts (avoids 'Duplicate name' on save)."""
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        prs.part.drop_rel(sldId.get(_R_ID))
        sldIdLst.remove(sldId)


RENDERERS = {
    "title": components.render_title,
    "thanks": components.render_thanks,
    "bullets": components.render_bullets,
    "card_row": components.render_card_row,
    "pill_flow": components.render_pill_flow,
    "callout": components.render_callout,
    "mapping": components.render_mapping,
    "figure": components.render_figure,
}


def _add_warning_slide(prs, message):
    slide = prs.slides.add_slide(prs.slide_layouts[ns.CONTENT_LAYOUT])
    ns.add_textbox(slide, ns.MARGIN_L, 1.0, ns.CONTENT_W, 0.5,
                   "BUILD WARNING", 18, "CF4F61", bold=True)
    ns.add_textbox(slide, ns.MARGIN_L, 1.8, ns.CONTENT_W, 3.0, message, 12, ns.INK)
    return slide


def build(outline, template_path, out_path):
    prs = load_template(template_path)
    clear_slides(prs)
    use_eyebrow = outline.get("meta", {}).get("use_eyebrow", True)
    for idx, page in enumerate(outline.get("slides", []), 1):
        page.setdefault("use_eyebrow", use_eyebrow)
        ptype = page.get("type")
        renderer = RENDERERS.get(ptype)
        if renderer is None:
            _add_warning_slide(prs, f"Slide {idx}: unknown type {ptype!r}; skipped.")
            continue
        try:
            renderer(prs, page)
        except Exception:
            _add_warning_slide(
                prs, f"Slide {idx} ({ptype}) failed to render:\n{traceback.format_exc()}")
    prs.save(out_path)
    return out_path


def main(argv=None):
    argv = argv or sys.argv[1:]
    if len(argv) != 2:
        print("usage: build_pptx.py <outline.json> <out.pptx>", file=sys.stderr)
        return 2
    outline_path, out_path = argv
    template = str(__import__("pathlib").Path(__file__).resolve().parents[1]
                   / "assets" / "northwestern_template.pptx")
    with open(outline_path, encoding="utf-8") as f:
        outline = json.load(f)
    path = build(outline, template, out_path)
    print(f"wrote {path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
