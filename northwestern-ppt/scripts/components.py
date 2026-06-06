"""One render_<type>(prs, page) -> slide per body type. All geometry/colors
come from nwstyle; these functions only place content."""
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import nwstyle as ns


def _content_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[ns.CONTENT_LAYOUT])


def _separator_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[ns.SEPARATOR_LAYOUT])


def header_band(slide, eyebrow, title, use_eyebrow=True):
    """Draw eyebrow + title + purple underline. Returns body_top (inches)."""
    if use_eyebrow and eyebrow:
        l, t, w, h = ns.EYEBROW_BOX
        ns.add_textbox(slide, l, t, w, h, eyebrow.upper(), 11, ns.PURPLE, bold=True)
    l, t, w, h = ns.TITLE_BOX
    ns.add_textbox(slide, l, t, w, h, title, 28, ns.PURPLE_DARK, bold=True)
    l, t, w, h = ns.UNDERLINE_BOX
    ns.add_rect(slide, l, t, w, h, ns.PURPLE)
    return ns.BODY_TOP


def render_title(prs, page):
    slide = _separator_slide(prs)
    ns.add_textbox(slide, 2.0, 2.0, 7.4, 1.2, page.get("title", ""), 34, ns.PURPLE_DARK, bold=True)
    if page.get("subtitle"):
        ns.add_textbox(slide, 2.0, 3.2, 7.4, 0.8, page["subtitle"], 18, ns.INK)
    if page.get("source"):
        ns.add_textbox(slide, 2.0, 4.6, 7.4, 0.4, page["source"], 12, ns.PURPLE)
    return slide


def render_thanks(prs, page):
    slide = _separator_slide(prs)
    ns.add_textbox(slide, 2.0, 2.4, 7.4, 1.0, page.get("title", "Thanks!"),
                   40, ns.PURPLE_DARK, bold=True)
    return slide


def render_bullets(prs, page):
    slide = _content_slide(prs)
    top = header_band(slide, page.get("eyebrow"), page.get("title", ""),
                      use_eyebrow=page.get("use_eyebrow", True))
    # fill the body region and vertically center so few bullets don't look sparse
    box = slide.shapes.add_textbox(ns.Inches(ns.MARGIN_L), ns.Inches(top + 0.1),
                                   ns.Inches(ns.CONTENT_W), ns.Inches(ns.BODY_BOTTOM - top - 0.1))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, item in enumerate(page.get("bullets", [])):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = ns.Pt(18)
        p.line_spacing = 1.15
        run = p.add_run()
        run.text = "•   " + item
        ns.style_run(run, 19, ns.INK)
    return slide


def render_card_row(prs, page):
    slide = _content_slide(prs)
    header_band(slide, page.get("eyebrow"), page.get("title", ""),
                use_eyebrow=page.get("use_eyebrow", True))
    cards = page.get("cards", [])[:3]
    n = max(len(cards), 1)
    gap = 0.3
    total_w = ns.CONTENT_W
    card_w = (total_w - gap * (n - 1)) / n
    # taller cards, vertically centered in the body region, so the slide reads full
    card_h = 2.5
    top = 2.05
    for i, card in enumerate(cards):
        line_hex, fill_hex = ns.ACCENTS[i % len(ns.ACCENTS)]
        left = ns.MARGIN_L + i * (card_w + gap)
        ns.add_rounded_rect(slide, left, top, card_w, card_h, fill_hex=fill_hex, line_hex=line_hex)
        ns.add_textbox(slide, left + 0.22, top + 0.22, card_w - 0.44, 0.4,
                       card.get("title", ""), 17, line_hex, bold=True)
        ns.add_textbox(slide, left + 0.22, top + 0.78, card_w - 0.44, card_h - 0.95,
                       card.get("desc", ""), 14, ns.INK)
    return slide


def render_pill_flow(prs, page):
    slide = _content_slide(prs)
    header_band(slide, page.get("eyebrow"), page.get("title", ""),
                use_eyebrow=page.get("use_eyebrow", True))
    pills = page.get("pills", [])
    n = max(len(pills), 1)
    pill_h = 0.85
    conn_w = 0.38
    top = 2.65
    total_w = ns.CONTENT_W
    pill_w = (total_w - conn_w * (n - 1)) / n
    x = ns.MARGIN_L
    cy = top + pill_h / 2
    line_hex = ns.PURPLE
    for i, label in enumerate(pills):
        shape = ns.add_rounded_rect(slide, x, top, pill_w, pill_h,
                                    fill_hex="F6F2FA", line_hex=line_hex)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = tf.paragraphs[0].add_run()
        run.text = label
        ns.style_run(run, 15, ns.PURPLE_DARK, bold=True)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        if i < n - 1:
            ns.add_connector(slide, x + pill_w, cy, x + pill_w + conn_w, cy, line_hex)
        x += pill_w + conn_w
    if page.get("note"):
        ns.add_textbox(slide, ns.MARGIN_L, top + pill_h + 0.5, ns.CONTENT_W, 0.6,
                       page["note"], 15, ns.INK)
    return slide


def render_callout(prs, page):
    slide = _content_slide(prs)
    header_band(slide, page.get("eyebrow"), page.get("title", ""),
                use_eyebrow=page.get("use_eyebrow", True))
    # large box filling the body region; big statement so the slide reads full
    box_top = 1.95
    box_h = 2.6
    ns.add_rounded_rect(slide, ns.MARGIN_L, box_top, ns.CONTENT_W, box_h,
                        fill_hex="F6F2FA", line_hex=ns.PURPLE)
    ns.add_textbox(slide, ns.MARGIN_L + 0.4, box_top + 0.3, ns.CONTENT_W - 0.8, 0.4,
                   page.get("label", ""), 14, ns.PURPLE, bold=True)
    ns.add_textbox(slide, ns.MARGIN_L + 0.4, box_top + 0.85, ns.CONTENT_W - 0.8, box_h - 1.1,
                   page.get("text", ""), 21, ns.INK)
    return slide


def render_mapping(prs, page):
    slide = _content_slide(prs)
    header_band(slide, page.get("eyebrow"), page.get("title", ""),
                use_eyebrow=page.get("use_eyebrow", True))
    rows = page.get("rows", [])
    col_gap = 0.25
    col_w = (ns.CONTENT_W - col_gap) / 2
    left_x = ns.MARGIN_L
    right_x = ns.MARGIN_L + col_w + col_gap
    top = 1.8
    head_h = 0.55
    row_gap = 0.14
    # size rows to spread across the body region so the slide reads full
    avail = ns.BODY_BOTTOM - (top + head_h + row_gap)
    n = max(len(rows), 1)
    row_h = max(0.55, min(0.95, (avail - row_gap * (n - 1)) / n))
    # header cells
    ns.add_rounded_rect(slide, left_x, top, col_w, head_h, fill_hex=ns.PURPLE, line_hex=None)
    ns.add_textbox(slide, left_x + 0.18, top + 0.1, col_w - 0.36, head_h - 0.2,
                   page.get("left_header", ""), 15, ns.WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    ns.add_rounded_rect(slide, right_x, top, col_w, head_h, fill_hex=ns.PURPLE, line_hex=None)
    ns.add_textbox(slide, right_x + 0.18, top + 0.1, col_w - 0.36, head_h - 0.2,
                   page.get("right_header", ""), 15, ns.WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    y = top + head_h + row_gap
    for left_text, right_text in rows:
        ns.add_rounded_rect(slide, left_x, y, col_w, row_h, fill_hex="F6F2FA", line_hex="E6E1EC")
        ns.add_textbox(slide, left_x + 0.18, y, col_w - 0.36, row_h,
                       left_text, 14, ns.INK, anchor=MSO_ANCHOR.MIDDLE)
        ns.add_rounded_rect(slide, right_x, y, col_w, row_h, fill_hex="F6F2FA", line_hex="E6E1EC")
        ns.add_textbox(slide, right_x + 0.18, y, col_w - 0.36, row_h,
                       right_text, 14, ns.INK, anchor=MSO_ANCHOR.MIDDLE)
        y += row_h + row_gap
    return slide


def render_table(prs, page):
    """A real PowerPoint table for tabular/numeric data (purple header row)."""
    from pptx.util import Inches
    slide = _content_slide(prs)
    header_band(slide, page.get("eyebrow"), page.get("title", ""),
                use_eyebrow=page.get("use_eyebrow", True))
    headers = page.get("headers", [])
    rows = page.get("rows", [])
    ncols = len(headers) or (len(rows[0]) if rows else 1)
    nrows = len(rows) + 1
    top = 1.85
    height = min(0.62 * nrows, ns.BODY_BOTTOM - top)
    gfx = slide.shapes.add_table(nrows, ncols, Inches(ns.MARGIN_L), Inches(top),
                                 Inches(ns.CONTENT_W), Inches(height))
    table = gfx.table
    table.first_row = False        # we style the header ourselves
    table.horz_banding = False
    for c, htext in enumerate(headers):
        _set_cell(table.cell(0, c), htext, 15, ns.WHITE, ns.PURPLE, bold=True)
    for r, row in enumerate(rows, start=1):
        fill = "FFFFFF" if r % 2 else "F6F2FA"
        for c in range(ncols):
            val = row[c] if c < len(row) else ""
            _set_cell(table.cell(r, c), str(val), 14, ns.INK, fill)
    return slide


def _set_cell(cell, text, size_pt, color_hex, fill_hex, bold=False):
    cell.fill.solid()
    cell.fill.fore_color.rgb = ns.hex_to_rgb(fill_hex)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_top = ns.Pt(3)
    cell.margin_bottom = ns.Pt(3)
    cell.margin_left = ns.Inches(0.15)
    p = cell.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    ns.style_run(run, size_pt, color_hex, bold=bold)
    return cell


def render_figure(prs, page):
    slide = _content_slide(prs)
    header_band(slide, page.get("eyebrow"), page.get("title", ""),
                use_eyebrow=page.get("use_eyebrow", True))
    has_cards = bool(page.get("cards"))
    img_left = ns.MARGIN_L
    img_top = 1.6
    img_max_w = (ns.CONTENT_W - 0.4) * (0.62 if has_cards else 1.0)
    img_max_h = 3.2
    # top-align the trimmed image; horizontally center it when it has the whole
    # width to itself so wide/narrow figures sit balanced under the title
    _pic, img_h = _add_scaled_picture(slide, page["image"], img_left, img_top,
                                      img_max_w, img_max_h, center=not has_cards)
    if page.get("caption"):
        ns.add_textbox(slide, img_left, img_top + img_h + 0.12, img_max_w, 0.4,
                       page["caption"], 12, ns.INK)
    if has_cards:
        card_left = ns.MARGIN_L + img_max_w + 0.4
        card_w = ns.CONTENT_W - img_max_w - 0.4
        y = img_top
        for i, card in enumerate(page["cards"][:3]):
            line_hex, fill_hex = ns.ACCENTS[i % len(ns.ACCENTS)]
            ch = 0.95
            ns.add_rounded_rect(slide, card_left, y, card_w, ch, fill_hex=fill_hex, line_hex=line_hex)
            ns.add_textbox(slide, card_left + 0.12, y + 0.1, card_w - 0.24, 0.3,
                           card.get("title", ""), 13, line_hex, bold=True)
            ns.add_textbox(slide, card_left + 0.12, y + 0.42, card_w - 0.24, ch - 0.5,
                           card.get("desc", ""), 11, ns.INK)
            y += ch + 0.2
    return slide


def _add_scaled_picture(slide, image_path, left, top, max_w, max_h, center=False):
    """Place an image scaled to fit within (max_w, max_h) inches, preserving aspect.

    White margins are auto-trimmed first (figures cropped from a PDF often carry
    whitespace). When center=True the scaled image is centered horizontally within
    max_w. Returns (picture_shape, placed_height_inches)."""
    import io
    from PIL import Image, ImageChops
    src = image_path
    try:
        im = Image.open(image_path).convert("RGB")
        bbox = ImageChops.difference(im, Image.new("RGB", im.size, (255, 255, 255))).getbbox()
        if bbox:
            im = im.crop(bbox)
        iw, ih = im.size
        buf = io.BytesIO()
        im.save(buf, format="PNG")
        buf.seek(0)
        src = buf
    except Exception:
        iw, ih = 4, 3
    aspect = iw / ih if ih else 4 / 3
    w = max_w
    h = w / aspect
    if h > max_h:
        h = max_h
        w = h * aspect
    if center:
        left = left + (max_w - w) / 2
    pic = slide.shapes.add_picture(src, ns.Inches(left), ns.Inches(top),
                                   ns.Inches(w), ns.Inches(h))
    return pic, h
