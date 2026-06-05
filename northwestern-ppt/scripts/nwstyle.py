"""Central style constants and low-level pptx shape/text helpers.

Every color/size/geometry value the deck uses lives here so the look stays
stable and editable in one place.
"""
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- template layout indices (verified against northwestern_template.pptx) ---
SEPARATOR_LAYOUT = 0   # left purple block bg -> title + section dividers
CONTENT_LAYOUT = 1     # purple footer bar + "Northwestern" wordmark -> body pages

# --- typography ---
FONT = "Times New Roman"

# --- palette ---
PURPLE = "542A84"
PURPLE_DARK = "3A1C5C"
INK = "262626"
WHITE = "FFFFFF"
ACCENTS = [
    ("4574B8", "EFF4FC"),  # blue
    ("3D916F", "EEF8F3"),  # green
    ("CE973E", "FCF7ED"),  # gold
    ("CF4F61", "FBEFF1"),  # coral
    ("542A84", "F6F2FA"),  # purple
]

# --- header band + body region geometry (inches) ---
EYEBROW_BOX = (0.55, 0.28, 8.9, 0.22)
TITLE_BOX = (0.55, 0.50, 9.0, 0.6)
UNDERLINE_BOX = (0.55, 1.12, 1.15, 0.045)
BODY_TOP = 1.45
BODY_BOTTOM = 5.05
MARGIN_L = 0.55
CONTENT_W = 8.9


def hex_to_rgb(h):
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def style_run(run, size_pt, color_hex, bold=False, font=FONT):
    run.font.name = font
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = hex_to_rgb(color_hex)


def add_textbox(slide, left, top, width, height, text, size_pt, color_hex,
                bold=False, align=PP_ALIGN.LEFT, font=FONT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    style_run(run, size_pt, color_hex, bold=bold, font=font)
    return box


def add_rounded_rect(slide, left, top, width, height, fill_hex=None,
                     line_hex=None, line_w_pt=1.0):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(height))
    shape.shadow.inherit = False
    if fill_hex is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(fill_hex)
    if line_hex is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = hex_to_rgb(line_hex)
        shape.line.width = Pt(line_w_pt)
    return shape


def add_rect(slide, left, top, width, height, fill_hex):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.shadow.inherit = False
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(fill_hex)
    shape.line.fill.background()
    return shape


def add_connector(slide, x1, y1, x2, y2, color_hex, w_pt=1.25):
    from pptx.enum.shapes import MSO_CONNECTOR
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = hex_to_rgb(color_hex)
    conn.line.width = Pt(w_pt)
    return conn
