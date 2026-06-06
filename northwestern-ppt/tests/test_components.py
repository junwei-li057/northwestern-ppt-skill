from pptx import Presentation
import build_pptx
import components
import nwstyle as ns


def fresh(template_path):
    prs = build_pptx.load_template(template_path)
    build_pptx.clear_slides(prs)
    return prs


def texts(slide):
    out = []
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            out.append(sh.text_frame.text.strip())
    return out


def test_title_slide(template_path):
    prs = fresh(template_path)
    page = {"type": "title", "title": "My Talk", "subtitle": "A subtitle", "source": "arXiv:1234"}
    slide = components.render_title(prs, page)
    joined = " | ".join(texts(slide))
    assert "My Talk" in joined and "A subtitle" in joined and "arXiv:1234" in joined
    assert slide.slide_layout == prs.slide_layouts[ns.SEPARATOR_LAYOUT]


def test_thanks_slide_default_text(template_path):
    prs = fresh(template_path)
    slide = components.render_thanks(prs, {"type": "thanks"})
    assert "Thanks!" in " | ".join(texts(slide))


def test_header_band_returns_body_top(template_path):
    prs = fresh(template_path)
    slide = prs.slides.add_slide(prs.slide_layouts[ns.CONTENT_LAYOUT])
    body_top = components.header_band(slide, "INTRO", "A title", use_eyebrow=True)
    assert abs(body_top - ns.BODY_TOP) < 1e-6
    assert "INTRO" in " | ".join(texts(slide))
    assert "A title" in " | ".join(texts(slide))


def test_bullets(template_path):
    prs = fresh(template_path)
    page = {"type": "bullets", "eyebrow": "METHOD", "title": "How it works",
            "bullets": ["First point", "Second point", "Third point"]}
    slide = components.render_bullets(prs, page)
    body = [t for t in texts(slide) if "point" in t.lower()]
    assert len(body) == 1                      # one body textbox holding all bullets
    for b in page["bullets"]:
        assert b in body[0]


def test_card_row_three_cards_rotating_accents(template_path):
    prs = fresh(template_path)
    page = {"type": "card_row", "eyebrow": "METRICS", "title": "Three scores",
            "cards": [{"title": "A", "desc": "alpha"},
                      {"title": "B", "desc": "beta"},
                      {"title": "C", "desc": "gamma"}]}
    slide = components.render_card_row(prs, page)
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    cards = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    fills = {c.fill.fore_color.rgb for c in cards}
    # the three card fills must all be present (other shapes like the underline
    # bar may also be filled, so check subset rather than exact equality)
    assert {ns.hex_to_rgb(ns.ACCENTS[i][1]) for i in range(3)} <= fills
    joined = " | ".join(texts(slide))
    for word in ["A", "alpha", "B", "beta", "C", "gamma"]:
        assert word in joined


def test_pill_flow(template_path):
    prs = fresh(template_path)
    page = {"type": "pill_flow", "eyebrow": "PIPELINE", "title": "Steps",
            "pills": ["Ingest", "Bin", "Sample", "Score"], "note": "left to right"}
    slide = components.render_pill_flow(prs, page)
    joined = " | ".join(texts(slide))
    for p in page["pills"]:
        assert p in joined
    assert "left to right" in joined
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    pills = [s for s in slide.shapes
             if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and s.has_text_frame
             and s.text_frame.text.strip() in page["pills"]]
    assert len(pills) == 4


def test_callout(template_path):
    prs = fresh(template_path)
    page = {"type": "callout", "eyebrow": "CLAIM", "title": "The point",
            "label": "Core claim", "text": "Process matters more than outcome."}
    slide = components.render_callout(prs, page)
    joined = " | ".join(texts(slide))
    assert "Core claim" in joined
    assert "Process matters more than outcome." in joined
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    boxes = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    assert any(b.fill.fore_color.rgb == ns.hex_to_rgb("F6F2FA") for b in boxes)


def test_mapping(template_path):
    prs = fresh(template_path)
    page = {"type": "mapping", "eyebrow": "DIAGNOSIS", "title": "Behavior to mode",
            "left_header": "Observed behavior", "right_header": "Attribution",
            "rows": [["No verification", "Task verification"],
                     ["Ignored input", "Misalignment"]]}
    slide = components.render_mapping(prs, page)
    joined = " | ".join(texts(slide))
    for cell in ["Observed behavior", "Attribution", "No verification",
                 "Task verification", "Ignored input", "Misalignment"]:
        assert cell in joined
