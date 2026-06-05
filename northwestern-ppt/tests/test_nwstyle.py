from pptx import Presentation
from pptx.util import Inches
import nwstyle


def _blank_slide():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    return prs, prs.slides.add_slide(prs.slide_layouts[6])  # 6 = blank in default template


def test_hex_to_rgb():
    rgb = nwstyle.hex_to_rgb("542A84")
    assert (rgb[0], rgb[1], rgb[2]) == (0x54, 0x2A, 0x84)


def test_add_textbox_sets_text_font_color():
    prs, slide = _blank_slide()
    shape = nwstyle.add_textbox(slide, 0.5, 0.5, 4, 0.5, "Hello", 18, "262626", bold=True)
    run = shape.text_frame.paragraphs[0].runs[0]
    assert run.text == "Hello"
    assert run.font.name == "Times New Roman"
    assert run.font.size.pt == 18
    assert run.font.bold is True
    assert run.font.color.rgb == nwstyle.hex_to_rgb("262626")


def test_add_rounded_rect_fill_and_line():
    prs, slide = _blank_slide()
    shape = nwstyle.add_rounded_rect(slide, 1, 1, 2, 1, fill_hex="EFF4FC", line_hex="4574B8")
    assert shape.fill.fore_color.rgb == nwstyle.hex_to_rgb("EFF4FC")
    assert shape.line.color.rgb == nwstyle.hex_to_rgb("4574B8")


def test_accents_has_five_pairs():
    assert len(nwstyle.ACCENTS) == 5
    for line_hex, fill_hex in nwstyle.ACCENTS:
        assert len(line_hex) == 6 and len(fill_hex) == 6
