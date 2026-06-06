import json
import warnings
from pptx import Presentation
import build_pptx


def test_load_template_has_three_sample_slides(template_path):
    prs = build_pptx.load_template(template_path)
    assert len(prs.slides._sldIdLst) == 3


def test_clear_slides_removes_all_without_duplicate_warning(template_path):
    prs = build_pptx.load_template(template_path)
    build_pptx.clear_slides(prs)
    assert len(prs.slides._sldIdLst) == 0


def test_rebuild_saves_clean_file(template_path, out_dir):
    prs = build_pptx.load_template(template_path)
    build_pptx.clear_slides(prs)
    prs.slides.add_slide(prs.slide_layouts[build_pptx.ns.CONTENT_LAYOUT])
    out = out_dir / "clean.pptx"
    with warnings.catch_warnings():
        warnings.simplefilter("error")     # duplicate-name -> error
        prs.save(str(out))
    reopened = Presentation(str(out))
    assert len(reopened.slides._sldIdLst) == 1


def test_build_full_outline(template_path, tmp_path):
    outline = {
        "meta": {"title": "T", "use_eyebrow": True},
        "slides": [
            {"type": "title", "title": "Talk", "subtitle": "sub", "source": "src"},
            {"type": "bullets", "eyebrow": "A", "title": "b", "bullets": ["x", "y"]},
            {"type": "callout", "eyebrow": "B", "title": "c", "label": "Claim", "text": "z"},
            {"type": "thanks"},
        ],
    }
    out = tmp_path / "deck.pptx"
    path = build_pptx.build(outline, template_path, str(out))
    prs = Presentation(path)
    assert len(prs.slides._sldIdLst) == 4


def test_build_unknown_type_adds_warning_slide(template_path, tmp_path):
    outline = {"meta": {}, "slides": [{"type": "nonsense", "title": "x"}]}
    out = tmp_path / "warn.pptx"
    build_pptx.build(outline, template_path, str(out))
    prs = Presentation(str(out))
    texts = []
    for slide in prs.slides:
        for sh in slide.shapes:
            if sh.has_text_frame:
                texts.append(sh.text_frame.text)
    assert any("WARNING" in t.upper() for t in texts)
